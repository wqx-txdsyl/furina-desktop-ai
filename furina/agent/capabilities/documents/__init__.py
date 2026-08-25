"""Documents capability（Phase 14E）—— 芙宁娜真正"写文件"。

优先 TXT / Markdown（真实 create/read/write/append/edit）；
DOCX / PPTX / XLSX 使用成熟库（python-docx / python-pptx / openpyxl），
保存后必须**重新打开验证**（reopen-verify）才算 verified；
PDF 本 Phase 只允许 read/extract 基础或标 unavailable（不引重型方案凑能力）。

输出必须返回真实 artifact path（进入 C7 Agent Task History）。
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import Any, Dict, Optional

from furina.agent.permission import Permission
from furina.agent.tool import BaseTool, ToolResult


def _resolve(path: str) -> Path:
    return Path(os.path.expanduser(path)).resolve()


# ================================================================ TXT / Markdown（原生）
class DocCreateTool(BaseTool):
    name = "doc.create"
    description = "创建 TXT/Markdown 文档并写入内容（L1，显式目标路径）"
    permission = Permission.L1_LOW_WRITE
    schema = {"type": "object", "properties": {"path": {"type": "string"},
                                               "content": {"type": "string"}},
              "required": ["path"]}

    def run(self, path: str, content: str = "") -> ToolResult:
        p = _resolve(path)
        if not p.suffix.lower() in (".txt", ".md", ".markdown"):
            return ToolResult(False, error=f"doc.create 只支持 TXT/Markdown: {p.suffix}", verified=False)
        if p.exists():
            return ToolResult(False, error=f"已存在（禁止 silent overwrite）: {p}", verified=False)
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(content or "", encoding="utf-8")
        except Exception as e:
            return ToolResult(False, error=str(e), verified=False)
        back = p.read_text(encoding="utf-8", errors="replace")
        return ToolResult(True, data={"path": str(p), "bytes": len(back.encode("utf-8"))},
                          verified=(back == (content or "")), note="创建并读回校验")


class DocReadTool(BaseTool):
    name = "doc.read"
    description = "读取 TXT/Markdown 文档内容（只读 L0）"
    permission = Permission.L0_READ
    schema = {"type": "object", "properties": {"path": {"type": "string"}},
              "required": ["path"]}

    def run(self, path: str) -> ToolResult:
        p = _resolve(path)
        if not p.is_file():
            return ToolResult(False, error=f"不是文件: {p}", verified=False)
        try:
            text = p.read_text(encoding="utf-8", errors="replace")
            return ToolResult(True, data={"path": str(p), "content": text[:10000],
                                          "bytes": len(text.encode("utf-8"))},
                              verified=True)
        except Exception as e:
            return ToolResult(False, error=str(e), verified=False)


class DocWriteTool(BaseTool):
    name = "doc.write"
    description = "覆盖写入 TXT/Markdown（显式目标；覆盖需 expected_old_hash 或 overwrite 语义，L1/L2）"
    permission = Permission.L1_LOW_WRITE
    schema = {"type": "object", "properties": {"path": {"type": "string"},
                                               "content": {"type": "string"},
                                               "expected_old_hash": {"type": "string"},
                                               "overwrite": {"type": "boolean"}},
              "required": ["path", "content"]}

    def run(self, path: str, content: str = "", expected_old_hash: str = "",
            overwrite: bool = False) -> ToolResult:
        import hashlib
        p = _resolve(path)
        if p.exists():
            if expected_old_hash:
                old = hashlib.sha256(p.read_bytes()).hexdigest()
                if old != expected_old_hash:
                    return ToolResult(False, error="expected_old_hash 不匹配，拒绝写入", verified=False)
            elif overwrite is False:
                return ToolResult(False, error="目标已存在且未显式 overwrite=True，拒绝覆盖",
                                  verified=False)
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(content or "", encoding="utf-8")
        except Exception as e:
            return ToolResult(False, error=str(e), verified=False)
        back = p.read_text(encoding="utf-8", errors="replace")
        return ToolResult(True, data={"path": str(p), "bytes": len(back.encode("utf-8"))},
                          verified=(back == (content or "")), note="写入并读回校验")


class DocAppendTool(BaseTool):
    name = "doc.append"
    description = "追加内容到 TXT/Markdown（L1）"
    permission = Permission.L1_LOW_WRITE
    schema = {"type": "object", "properties": {"path": {"type": "string"},
                                               "content": {"type": "string"}},
              "required": ["path", "content"]}

    def run(self, path: str, content: str = "") -> ToolResult:
        p = _resolve(path)
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            with open(p, "a", encoding="utf-8") as f:
                f.write(content or "")
        except Exception as e:
            return ToolResult(False, error=str(e), verified=False)
        back = p.read_text(encoding="utf-8", errors="replace")
        return ToolResult(True, data={"path": str(p)},
                          verified=(content in back), note="追加并校验")


class DocEditTool(BaseTool):
    name = "doc.edit"
    description = "在 TXT/Markdown 中替换子串（L1/L2）"
    permission = Permission.L1_LOW_WRITE
    schema = {"type": "object", "properties": {"path": {"type": "string"},
                                               "old": {"type": "string"},
                                               "new": {"type": "string"}},
              "required": ["path", "old", "new"]}

    def run(self, path: str, old: str, new: str) -> ToolResult:
        p = _resolve(path)
        if not p.is_file():
            return ToolResult(False, error=f"不是文件: {p}", verified=False)
        try:
            text = p.read_text(encoding="utf-8", errors="replace")
            if old not in text:
                return ToolResult(False, error="未找到要替换的内容", verified=False)
            p.write_text(text.replace(old, new, 1), encoding="utf-8")
        except Exception as e:
            return ToolResult(False, error=str(e), verified=False)
        back = p.read_text(encoding="utf-8", errors="replace")
        return ToolResult(True, data={"path": str(p)},
                          verified=(new in back and text != back),
                          note="编辑并校验")


# ================================================================ DOCX（python-docx + reopen-verify）
class DocxCreateTool(BaseTool):
    name = "docx.create"
    description = "创建 Word 文档：标题/段落/简单列表，保存后重新打开验证段落内容（L1）"
    permission = Permission.L1_LOW_WRITE
    schema = {"type": "object", "properties": {"path": {"type": "string"},
                                               "title": {"type": "string"},
                                               "paragraphs": {"type": "array"},
                                               "bullets": {"type": "array"}},
              "required": ["path"]}

    def run(self, path: str, title: str = "", paragraphs: Optional[list] = None,
            bullets: Optional[list] = None) -> ToolResult:
        try:
            from docx import Document
        except Exception as e:
            return ToolResult(False, error=f"python-docx 不可用: {e}", verified=False)
        p = _resolve(path)
        if not p.suffix.lower() == ".docx":
            return ToolResult(False, error="docx.create 需要 .docx 路径", verified=False)
        try:
            doc = Document()
            if title:
                doc.add_heading(title or "", level=0)
            for para in (paragraphs or []):
                doc.add_paragraph(str(para))
            for b in (bullets or []):
                doc.add_paragraph(str(b), style="List Bullet")
            p.parent.mkdir(parents=True, exist_ok=True)
            doc.save(str(p))
        except Exception as e:
            return ToolResult(False, error=str(e), verified=False)
        # reopen-verify：重新打开验证段落内容
        try:
            back = Document(str(p))
            texts = [para.text for para in back.paragraphs]
            ok = (title in texts or not title) and all(t in texts for t in (paragraphs or []))
        except Exception as e:
            return ToolResult(False, error=f"重新打开验证失败: {e}", verified=False)
        return ToolResult(True, data={"path": str(p), "paragraphs": len(texts)},
                          verified=ok, note="保存并重新打开验证段落内容")


# ================================================================ PPTX（python-pptx + reopen-verify）
class PptxCreateTool(BaseTool):
    name = "pptx.create"
    description = "创建 PowerPoint：title/content slides，保存后重新打开验证 slide count（L1）"
    permission = Permission.L1_LOW_WRITE
    schema = {"type": "object", "properties": {"path": {"type": "string"},
                                               "slides": {"type": "array"}},
              "required": ["path"]}

    def run(self, path: str, slides: Optional[list] = None) -> ToolResult:
        try:
            from pptx import Presentation
        except Exception as e:
            return ToolResult(False, error=f"python-pptx 不可用: {e}", verified=False)
        p = _resolve(path)
        if not p.suffix.lower() == ".pptx":
            return ToolResult(False, error="pptx.create 需要 .pptx 路径", verified=False)
        slides = slides or [{"title": "标题", "bullets": []}]
        try:
            prs = Presentation()
            for s in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
                slide.shapes.title.text = str(s.get("title", ""))
                body = slide.placeholders[1]
                body.text_frame.text = ""
                for i, b in enumerate(s.get("bullets", []) or []):
                    if i == 0:
                        body.text_frame.text = str(b)
                    else:
                        body.text_frame.add_paragraph().text = str(b)
            p.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(p))
        except Exception as e:
            return ToolResult(False, error=str(e), verified=False)
        # reopen-verify：slide count
        try:
            back = Presentation(str(p))
            n = len(back.slides)
            ok = n == len(slides)
        except Exception as e:
            return ToolResult(False, error=f"重新打开验证失败: {e}", verified=False)
        return ToolResult(True, data={"path": str(p), "slides": n},
                          verified=ok, note="保存并重新打开验证 slide count")


# ================================================================ XLSX（openpyxl + reopen-verify）
class XlsxCreateTool(BaseTool):
    name = "xlsx.create"
    description = "创建 Excel：写二维数据，保存后重新打开验证 cell values（L1）"
    permission = Permission.L1_LOW_WRITE
    schema = {"type": "object", "properties": {"path": {"type": "string"},
                                               "rows": {"type": "array"},
                                               "sheet": {"type": "string"}},
              "required": ["path"]}

    def run(self, path: str, rows: Optional[list] = None, sheet: str = "Sheet1") -> ToolResult:
        try:
            from openpyxl import Workbook, load_workbook
        except Exception as e:
            return ToolResult(False, error=f"openpyxl 不可用: {e}", verified=False)
        p = _resolve(path)
        if not p.suffix.lower() in (".xlsx", ".xlsm"):
            return ToolResult(False, error="xlsx.create 需要 .xlsx 路径", verified=False)
        rows = rows or []
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = (sheet or "Sheet1")[:31]
            for r in rows:
                ws.append([c for c in (r or [])])
            p.parent.mkdir(parents=True, exist_ok=True)
            wb.save(str(p))
        except Exception as e:
            return ToolResult(False, error=str(e), verified=False)
        # reopen-verify：cell values
        try:
            wb2 = load_workbook(str(p), read_only=True)
            ws2 = wb2.active
            vals = [list(row) for row in ws2.iter_rows(values_only=True)]
            ok = vals == [list(r) for r in rows]
            wb2.close()
        except Exception as e:
            return ToolResult(False, error=f"重新打开验证失败: {e}", verified=False)
        return ToolResult(True, data={"path": str(p), "cells": len(rows)},
                          verified=ok, note="保存并重新打开验证 cell values")


ALL_DOC_TOOLS = [DocCreateTool, DocReadTool, DocWriteTool, DocAppendTool, DocEditTool,
                 DocxCreateTool, PptxCreateTool, XlsxCreateTool]
